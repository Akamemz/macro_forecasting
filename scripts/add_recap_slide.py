"""
Insert a "Project Recap" slide at position 0 in progress_week1.pptx.
Matches the dark-navy header + body style of the existing slides.

Usage:
    python scripts/add_recap_slide.py
"""

import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
PPTX = ROOT / "research" / "ppt" / "progress_week1.pptx"

NAVY   = RGBColor(0x1a, 0x1a, 0x2e)
DARK2  = RGBColor(0x16, 0x21, 0x3e)   # accent panels
ACCENT = RGBColor(0x2d, 0x4a, 0x7a)   # blue accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)
GREY   = RGBColor(0x88, 0x88, 0x88)

EMU = 914400  # 1 inch in EMU



def add_rect(slide, l, t, w, h, color: RGBColor):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(l * EMU), Emu(t * EMU), Emu(w * EMU), Emu(h * EMU),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, l, t, w, h, text, size=18, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(
        Emu(l * EMU), Emu(t * EMU), Emu(w * EMU), Emu(h * EMU)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_bullet_box(slide, l, t, w, h, lines, size=11, header=None, header_size=12):
    """Add a panel with an optional bold header and bullet lines."""
    txb = slide.shapes.add_textbox(
        Emu(l * EMU), Emu(t * EMU), Emu(w * EMU), Emu(h * EMU)
    )
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    if header:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = header
        r.font.size  = Pt(header_size)
        r.font.bold  = True
        r.font.color.rgb = WHITE

    for line in lines:
        p = tf.paragraphs[0] if (first and not header) else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size  = Pt(size)
        r.font.bold  = False
        r.font.color.rgb = LIGHT

    return txb


def move_slide_to_front(prs, index):
    """Move slide at `index` to position 0."""
    from pptx.oxml.ns import qn
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    elem = slides[index]
    xml_slides.remove(elem)
    xml_slides.insert(0, elem)


def build_recap_slide(prs):
    blank_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(blank_layout)

    W, H = 13.33, 7.50   # slide dimensions in inches

    # ── dark background ───────────────────────────────────────────────────────
    bg = add_rect(slide, 0, 0, W, H, NAVY)

    # ── top header bar (1.25 in tall, matches other slides) ───────────────────
    add_rect(slide, 0, 0, W, 1.25, DARK2)

    # ── title ─────────────────────────────────────────────────────────────────
    add_textbox(slide, 0.40, 0.12, 10.0, 0.60,
                "Project Recap", size=28, bold=True, color=WHITE)
    add_textbox(slide, 0.40, 0.72, 12.0, 0.40,
                "Macro Forecasting with IMF Text  ·  STAT 8240  ·  Kennesaw State University  ·  April 2025",
                size=10, color=RGBColor(0xAA, 0xBB, 0xCC))

    # ── research question banner ──────────────────────────────────────────────
    add_rect(slide, 0.40, 1.45, W - 0.80, 0.65, ACCENT)
    add_textbox(slide, 0.55, 1.50, W - 1.10, 0.55,
                "Research question:  Does adding IMF Article IV staff report text improve macroeconomic "
                "forecasts over a numerical-only baseline?",
                size=12, bold=False, italic=True, color=WHITE)

    # ── two-column body ───────────────────────────────────────────────────────
    COL_L  = 0.40          # left col x
    COL_R  = 6.90          # right col x
    CW     = 6.20          # col width
    TOP    = 2.30          # body top

    # left: Data & Countries
    add_rect(slide, COL_L, TOP, CW, 0.38, ACCENT)
    add_textbox(slide, COL_L + 0.12, TOP + 0.05, CW - 0.2, 0.30,
                "DATA STREAMS", size=10, bold=True, color=WHITE)

    add_bullet_box(
        slide, COL_L + 0.10, TOP + 0.45, CW - 0.20, 1.70,
        lines=[
            "① World Bank WDI  — 7 annual macro indicators per country",
            "    GDP growth, inflation, current account, fiscal balance,",
            "    remittances, unemployment, government debt  (2005–2023)",
            "",
            "② IMF Article IV staff reports  — qualitative economic assessment",
            "    Published 3–8 months after the reference year",
            "    240 real texts extracted  ·  57 % coverage across 22 countries",
        ],
        size=10.5,
    )

    # left: Countries box
    add_rect(slide, COL_L, TOP + 2.30, CW, 0.38, ACCENT)
    add_textbox(slide, COL_L + 0.12, TOP + 2.35, CW - 0.2, 0.30,
                "COUNTRIES  (22 total)", size=10, bold=True, color=WHITE)
    add_bullet_box(
        slide, COL_L + 0.10, TOP + 2.75, CW - 0.20, 1.10,
        lines=[
            "Central Asia: KAZ  KGZ  TJK  TKM  UZB  AZE  ARM  GEO  MDA  BLR  UKR  MNG",
            "Eastern Europe: ALB  BIH  MKD  SRB  TUR  ROU  BGR  HRV",
            "Anchors: USA  JPN",
        ],
        size=10.5,
    )

    # right: Models
    add_rect(slide, COL_R, TOP, CW, 0.38, ACCENT)
    add_textbox(slide, COL_R + 0.12, TOP + 0.05, CW - 0.2, 0.30,
                "MODELS", size=10, bold=True, color=WHITE)
    add_bullet_box(
        slide, COL_R + 0.10, TOP + 0.45, CW - 0.20, 1.80,
        lines=[
            "DLinear  —  trend decomposition + linear projection",
            "    Unimodal baseline (numbers only)",
            "",
            "NumericalGRU  —  GRU over WDI time series",
            "    Unimodal ablation: temporal dynamics without text",
            "",
            "GR-Add MMF  —  GRU + RecAvg TTF + gated text fusion",
            "    Full multimodal model (numbers + Article IV text)",
        ],
        size=10.5,
    )

    # right: Evaluation
    add_rect(slide, COL_R, TOP + 2.30, CW, 0.38, ACCENT)
    add_textbox(slide, COL_R + 0.12, TOP + 2.35, CW - 0.2, 0.30,
                "EVALUATION", size=10, bold=True, color=WHITE)
    add_bullet_box(
        slide, COL_R + 0.10, TOP + 2.75, CW - 0.20, 1.10,
        lines=[
            "Targets: next-year GDP growth  ·  next-year CPI inflation",
            "Train 2005–2019  |  Val 2020–2021  |  Test 2022–2023",
            "Metric: MSE on held-out test years",
        ],
        size=10.5,
    )

    # ── footer ────────────────────────────────────────────────────────────────
    add_textbox(slide, 0.40, 7.10, W - 0.80, 0.28,
                "STAT 8240 Data Mining II  ·  Kennesaw State University  ·  Track B Application of Existing Methods to New Datasets",
                size=8, color=GREY, align=PP_ALIGN.CENTER)

    return slide


def main():
    prs = Presentation(str(PPTX))
    print(f"Slides before: {len(prs.slides)}")

    build_recap_slide(prs)

    # Move the newly added slide (last position) to position 0
    move_slide_to_front(prs, len(prs.slides) - 1)

    prs.save(str(PPTX))
    print(f"Slides after:  {len(prs.slides)}")
    print(f"Saved: {PPTX}")


if __name__ == "__main__":
    main()

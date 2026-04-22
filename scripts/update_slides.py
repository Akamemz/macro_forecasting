"""Update slide 1 of progress_week1.pptx to reflect Week 2 real-data results."""
import copy
from pptx import Presentation
from pptx.util import Pt
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
PPTX = ROOT / "research" / "ppt" / "progress_week1.pptx"

prs = Presentation(str(PPTX))
slide = prs.slides[0]

# ── helper: replace all text in a shape while keeping first-run formatting ────
def set_text(shape, new_text: str):
    tf = shape.text_frame
    # copy formatting from the first run of the first paragraph
    first_para = tf.paragraphs[0]
    first_run  = first_para.runs[0] if first_para.runs else None

    # clear all paragraphs
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""

    # write the new text into the first paragraph's first run
    if first_run:
        first_run.text = new_text
    else:
        first_para.text = new_text

def set_multiline(shape, lines: list[str]):
    """Replace text frame content with multiple lines, one paragraph each."""
    from pptx.oxml.ns import qn
    import copy
    tf = shape.text_frame

    # grab format template from first paragraph/run
    tmpl_para = tf.paragraphs[0]
    tmpl_run  = tmpl_para.runs[0] if tmpl_para.runs else None

    # clear existing paragraphs text
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""

    # set first paragraph
    if tmpl_run:
        tmpl_run.text = lines[0]
    else:
        tmpl_para.text = lines[0]

    # add extra paragraphs for remaining lines
    txBody = tf._txBody
    for line in lines[1:]:
        new_para = copy.deepcopy(tmpl_para._p)
        # clear runs in new para and set text
        for r in new_para.findall(qn('a:r')):
            for t in r.findall(qn('a:t')):
                t.text = ""
        runs = new_para.findall(qn('a:r'))
        if runs:
            runs[0].find(qn('a:t')).text = line
        txBody.append(new_para)

# ── index shapes by name ──────────────────────────────────────────────────────
shapes = {s.name: s for s in slide.shapes if s.has_text_frame}

# ── 1. Title ──────────────────────────────────────────────────────────────────
set_text(shapes["TextBox 2"], "Week 2 Progress: Real Data Pipeline Complete")

# ── 2. Subtitle ───────────────────────────────────────────────────────────────
set_text(shapes["TextBox 3"], "Macro Forecasting with IMF Text  ·  STAT 8240  ·  April 2025")

# ── 3. Dummy Data Generation → Real Data Acquisition ─────────────────────────
set_text(shapes["TextBox 10"], "Real Data Acquisition")
set_text(shapes["TextBox 11"],
    "World Bank WDI API: 7 indicators × 22 countries × 2005–2023. "
    "IMF Article IV reports: Coveo search API → 240 real staff report texts extracted "
    "(57% coverage). PDFs streamed in memory — zero disk storage for raw PDFs.")

# ── 4. PDF Extraction → update to reflect streaming approach ──────────────────
set_text(shapes["TextBox 8"],
    "PyMuPDF extracts pages 1–6 (Executive Summary + Staff Appraisal) directly from "
    "in-memory PDF bytes — no files written to disk. Coveo API used to locate reports "
    "per country; Playwright browser automates the IMF download button.")

# ── 5. DATA AT A GLANCE stats ─────────────────────────────────────────────────
set_text(shapes["TextBox 30"], "240")
set_text(shapes["TextBox 31"], "real Article IV texts")
set_text(shapes["TextBox 33"], "57%")
set_text(shapes["TextBox 34"], "text coverage (22 countries)")

# ── 6. Footer: update Next steps ─────────────────────────────────────────────
set_text(shapes["TextBox 38"],
    "Next: run BERT encoding on real Article IV texts  ·  full 100-epoch experiment matrix  ·  ablation analysis")

# ── Speaker notes ─────────────────────────────────────────────────────────────
def set_notes(slide, text: str):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

NOTES_SLIDE1 = (
    "Week 2 delivered the full real-data pipeline on top of the dummy scaffold built in Week 1. "
    "Three major pieces of work were completed.\n\n"

    "First, World Bank WDI data: we hit the World Bank API directly to pull all seven macro indicators "
    "(GDP growth, inflation, fiscal balance, current account, remittances, unemployment, government debt) "
    "for all 22 countries across 2005–2023. The existing CSV already covered most ECA countries; "
    "we appended USA, Japan, and Mongolia which were missing.\n\n"

    "Second, IMF Article IV text acquisition: rather than downloading thousands of PDFs, we reverse-engineered "
    "the IMF website's Coveo search API. This let us query per country — roughly 30–100 hits each — "
    "instead of paginating through 7,335 total results. For each matched report, a Playwright browser "
    "clicks the IMF download button, intercepts the PDF bytes in memory, PyMuPDF extracts pages 1–6, "
    "and only the cleaned text is saved. No PDFs are stored on disk — this keeps the repo lightweight.\n\n"

    "Third, coverage: we achieved 240 real Article IV texts across 22 countries, a 57% fill rate. "
    "The gaps are genuine — Turkmenistan never publishes its reports publicly, Belarus and Ukraine "
    "have politically suspended consultations in recent years, and many Central Asian countries are on "
    "IMF programs where Article IV is folded into biennial program reviews rather than annual standalone reports. "
    "The pipeline's missingness mask handles these gaps explicitly, so missing text years simply default to "
    "the numerical-only pathway.\n\n"

    "Models are already implemented and were tested end-to-end on dummy data in Week 1. "
    "Week 3 begins BERT encoding of the real texts and launches the full experiment matrix."
)

NOTES_SLIDE2 = (
    "This slide walks through the four-step IMM-TSF pipeline as applied to our macro forecasting problem. "
    "Each step maps directly to a module in the codebase.\n\n"

    "Step 1 — Pre-alignment: for each (country, query year) pair we extract a 10-year lookback window "
    "of WDI data, producing a (10 × 7) value matrix. Missing observations are flagged in a binary mask "
    "rather than imputed, and calendar years are normalised to [0, 1] so the model sees relative time, "
    "not raw years.\n\n"

    "Step 2 — Text encoding: each Article IV staff report summary is passed through frozen BERT-base-uncased "
    "to produce a 768-dimensional embedding. Embeddings are precomputed once and stored as .npy files "
    "so training never re-runs the language model. For country-years with no report, the embedding is a "
    "zero vector and the missingness mask signals its absence.\n\n"

    "Step 3 — RecAvg TTF: the Temporally-Tagged Fusion module applies a Gaussian kernel (σ = 1.0 year) "
    "to weight past embeddings by their distance from the query year. This naturally handles the 3–8 month "
    "IMF publication lag — a 2022 Article IV report published in mid-2023 contributes less weight to a "
    "2023 forecast query than a report published on time would.\n\n"

    "Step 4 — GR-Add MMF: the Gated Residual Multimodal Fusion layer combines the GRU hidden state "
    "(numerical stream) with the projected text vector. A sigmoid gate decides how much of the text delta "
    "to add. If text is missing the gate collapses and the model falls back to a pure GRU prediction.\n\n"

    "Evaluation uses a strict chronological split: train on 2005–2019, validate on 2020–2021, test on 2022–2023. "
    "Primary metric is MSE on the held-out test years for GDP growth and inflation. "
    "Secondary analysis stratifies the MSE gain by text missingness severity to test whether gappier "
    "countries benefit more from the multimodal signal when text is available."
)

set_notes(prs.slides[0], NOTES_SLIDE1)
set_notes(prs.slides[1], NOTES_SLIDE2)

prs.save(str(PPTX))
print(f"Saved: {PPTX}")
